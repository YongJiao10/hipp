#!/opt/miniconda3/envs/py314/bin/python
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(0, 0, 0)
FG = RGBColor(255, 255, 255)
MUTED = RGBColor(170, 170, 170)
ACCENT = RGBColor(90, 90, 90)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str, size: int, bold: bool = False, color: RGBColor = FG, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Arial"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color


def add_image_fit(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    with Image.open(path) as img:
        img_w, img_h = img.size
    scale = min(width / img_w, height / img_h)
    draw_w = img_w * scale
    draw_h = img_h * scale
    x = left + (width - draw_w) / 2.0
    y = top + (height - draw_h) / 2.0
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(draw_w), height=Inches(draw_h))


def add_subject_slide(prs: Presentation, subject: str, structural_path: Path, gradient_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_textbox(slide, 0.45, 0.25, 4.0, 0.45, f"sub-{subject}", 24, bold=True)
    add_textbox(slide, 0.45, 0.68, 3.0, 0.28, "Locked native views", 11, color=MUTED)

    left_x = 0.45
    gap_x = 0.32
    panel_w = 6.05
    img_top = 1.18
    img_h = 5.55
    right_x = left_x + panel_w + gap_x

    add_textbox(slide, left_x, 0.92, panel_w, 0.24, "Structural", 14, bold=True)
    add_textbox(slide, right_x, 0.92, panel_w, 0.24, "Gradient 1", 14, bold=True)

    add_image_fit(slide, structural_path, left_x, img_top, panel_w, img_h)
    add_image_fit(slide, gradient_path, right_x, img_top, panel_w, img_h)

    slide.shapes.add_shape(
        1,
        Inches(left_x),
        Inches(7.02),
        Inches(panel_w),
        Inches(0.015),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ACCENT
    slide.shapes[-1].line.fill.background()

    slide.shapes.add_shape(
        1,
        Inches(right_x),
        Inches(7.02),
        Inches(panel_w),
        Inches(0.015),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ACCENT
    slide.shapes[-1].line.fill.background()


def main() -> int:
    repo = Path("/Users/jy/Documents/HippoMaps")
    final_root = repo / "outputs" / "dense_corobl_batch" / "final_wb_locked"
    out_path = final_root / "HippoMaps_locked_views_summary.pptx"

    subjects = ["100610", "102311", "102816"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "HippoMaps Locked Views Summary"
    prs.core_properties.subject = "Formal structural and hippocampal FC gradient views"
    prs.core_properties.author = "OpenAI Codex"

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(title_slide)
    add_textbox(title_slide, 0.7, 1.15, 12.0, 0.8, "HippoMaps Locked Views Summary", 28, bold=True)
    add_textbox(title_slide, 0.7, 2.0, 10.5, 0.4, "Six finalized figures across three subjects", 15, color=MUTED)
    add_textbox(title_slide, 0.7, 2.55, 11.0, 0.4, "Each subject is summarized on one slide with structural and primary FC gradient views side by side.", 14, color=FG)
    add_textbox(title_slide, 0.7, 6.55, 5.0, 0.3, "Source: outputs/dense_corobl_batch/final_wb_locked", 11, color=MUTED)

    for subject in subjects:
        subject_dir = final_root / f"sub-{subject}"
        add_subject_slide(
            prs,
            subject,
            subject_dir / f"sub-{subject}_structural.png",
            subject_dir / f"sub-{subject}_gradient.png",
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
